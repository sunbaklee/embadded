from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "presentation_assets"
SHOTS = ASSETS / "screenshots"
CROPS = ASSETS / "crops"
OUTPUT = ROOT / "IoT_LoneCare_프로젝트_발표자료.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "navy": RGBColor(18, 60, 52),
    "green": RGBColor(22, 139, 103),
    "mint": RGBColor(84, 215, 165),
    "pale": RGBColor(241, 247, 244),
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(26, 36, 33),
    "muted": RGBColor(91, 110, 104),
    "line": RGBColor(215, 226, 221),
    "red": RGBColor(204, 57, 70),
    "red_pale": RGBColor(255, 235, 238),
    "amber": RGBColor(184, 118, 0),
    "amber_pale": RGBColor(255, 244, 218),
    "blue": RGBColor(62, 126, 209),
    "blue_pale": RGBColor(232, 241, 253),
    "gray": RGBColor(247, 249, 248),
}

FONT = "맑은 고딕"


def crop_image(source: Path, output: Path, box: tuple[int, int, int, int]) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(source) as image:
        image.crop(box).save(output, quality=94)


def prepare_crops() -> dict[str, Path]:
    paths = {
        "dashboard_top": CROPS / "dashboard_top.png",
        "dashboard_alerts": CROPS / "dashboard_alerts.png",
        "devices": CROPS / "devices.png",
        "docs": CROPS / "docs.png",
        "resolutions": CROPS / "resolutions.png",
    }
    crop_image(SHOTS / "dashboard.png", paths["dashboard_top"], (60, 20, 1380, 940))
    crop_image(
        SHOTS / "dashboard.png",
        paths["dashboard_alerts"],
        (80, 850, 1360, 2190),
    )
    crop_image(SHOTS / "devices.png", paths["devices"], (35, 30, 1405, 1510))
    crop_image(SHOTS / "docs.png", paths["docs"], (15, 40, 1425, 1700))
    crop_image(SHOTS / "resolutions.png", paths["resolutions"], (100, 25, 1340, 1280))
    return paths


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill,
    radius=True,
    line=None,
    line_width=1,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT,
    margin=0.04,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_rich_text(
    slide,
    runs,
    x,
    y,
    w,
    h,
    size=20,
    color=None,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, bold, run_color in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = run_color or color or COLORS["ink"]
    return box


def add_bullets(slide, items, x, y, w, h, size=18, color=None, gap=8):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {item}"
        paragraph.level = 0
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color or COLORS["ink"]
        paragraph.space_after = Pt(gap)
    return box


def add_title(slide, number: str, title: str, subtitle: str | None = None):
    add_text(
        slide,
        number,
        Inches(0.55),
        Inches(0.38),
        Inches(0.5),
        Inches(0.32),
        size=11,
        color=COLORS["green"],
        bold=True,
    )
    add_text(
        slide,
        title,
        Inches(0.55),
        Inches(0.67),
        Inches(11.9),
        Inches(0.55),
        size=28,
        color=COLORS["navy"],
        bold=True,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.58),
            Inches(1.23),
            Inches(11.8),
            Inches(0.35),
            size=12,
            color=COLORS["muted"],
        )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(1.58),
        Inches(12.2),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_footer(slide, page: int, source: str = "프로젝트 계획서 · README · 코드 분석"):
    add_text(
        slide,
        source,
        Inches(0.58),
        Inches(7.15),
        Inches(5.6),
        Inches(0.2),
        size=8,
        color=COLORS["muted"],
    )
    add_text(
        slide,
        f"{page:02d}",
        Inches(12.1),
        Inches(7.12),
        Inches(0.6),
        Inches(0.22),
        size=9,
        color=COLORS["green"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_picture_cover(slide, path: Path, x, y, w, h, line=True):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = w / h
    if image_ratio > frame_ratio:
        pic_h = h
        pic_w = h * image_ratio
        pic_x = x - (pic_w - w) / 2
        pic_y = y
    else:
        pic_w = w
        pic_h = w / image_ratio
        pic_x = x
        pic_y = y - (pic_h - h) / 2
    pic = slide.shapes.add_picture(str(path), pic_x, pic_y, pic_w, pic_h)
    overlay = add_rect(
        slide,
        x,
        y,
        w,
        h,
        RGBColor(255, 255, 255),
        radius=True,
        line=COLORS["line"] if line else RGBColor(255, 255, 255),
    )
    overlay.fill.background()
    overlay.line.color.rgb = COLORS["line"] if line else COLORS["white"]
    pic.element.addprevious(overlay.element)
    return pic


def add_arch_box(slide, title, detail, x, y, w, color, number):
    add_rect(
        slide,
        x,
        y,
        w,
        Inches(1.25),
        COLORS["white"],
        line=COLORS["line"],
    )
    add_rect(
        slide,
        x + Inches(0.18),
        y + Inches(0.2),
        Inches(0.42),
        Inches(0.42),
        color,
    )
    add_text(
        slide,
        number,
        x + Inches(0.18),
        y + Inches(0.2),
        Inches(0.42),
        Inches(0.42),
        size=14,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title,
        x + Inches(0.72),
        y + Inches(0.17),
        w - Inches(0.9),
        Inches(0.38),
        size=16,
        color=COLORS["navy"],
        bold=True,
    )
    add_text(
        slide,
        detail,
        x + Inches(0.22),
        y + Inches(0.68),
        w - Inches(0.44),
        Inches(0.42),
        size=10.5,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def connect(slide, x1, y1, x2, y2, color=None):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color or COLORS["green"]
    connector.line.width = Pt(2)
    connector.line.end_arrowhead = True
    return connector


def add_status_card(slide, label, value, note, x, y, color, pale):
    add_rect(
        slide,
        x,
        y,
        Inches(2.75),
        Inches(1.55),
        COLORS["white"],
        line=COLORS["line"],
    )
    add_rect(
        slide,
        x,
        y,
        Inches(0.07),
        Inches(1.55),
        color,
        radius=False,
    )
    add_rect(
        slide,
        x + Inches(2.12),
        y + Inches(0.77),
        Inches(0.48),
        Inches(0.48),
        pale,
    )
    add_text(
        slide,
        label,
        x + Inches(0.2),
        y + Inches(0.27),
        Inches(1.7),
        Inches(0.25),
        size=11,
        color=COLORS["muted"],
        bold=True,
    )
    add_text(
        slide,
        value,
        x + Inches(0.2),
        y + Inches(0.55),
        Inches(1.5),
        Inches(0.55),
        size=28,
        color=color,
        bold=True,
    )
    add_text(
        slide,
        note,
        x + Inches(0.2),
        y + Inches(1.15),
        Inches(2.0),
        Inches(0.25),
        size=10,
        color=COLORS["muted"],
    )


def create_presentation() -> Path:
    crops = prepare_crops()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["navy"], radius=False)
    add_rect(
        slide,
        Inches(8.55),
        Inches(-0.55),
        Inches(5.7),
        Inches(8.4),
        RGBColor(25, 91, 76),
        radius=False,
    )
    add_rect(
        slide,
        Inches(9.2),
        Inches(0.45),
        Inches(3.5),
        Inches(6.2),
        RGBColor(34, 111, 92),
        line=RGBColor(34, 111, 92),
    )
    add_text(
        slide,
        "EMBEDDED SYSTEM PROJECT",
        Inches(0.75),
        Inches(0.72),
        Inches(5.5),
        Inches(0.3),
        size=11,
        color=COLORS["mint"],
        bold=True,
    )
    add_text(
        slide,
        "IoT LoneCare",
        Inches(0.72),
        Inches(1.35),
        Inches(7.2),
        Inches(0.9),
        size=42,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "비접촉식 IoT 센서 기반\n고독사 예방 및 생활 안전 모니터링 시스템",
        Inches(0.76),
        Inches(2.35),
        Inches(6.8),
        Inches(1.2),
        size=24,
        color=RGBColor(225, 244, 237),
        bold=True,
    )
    add_text(
        slide,
        "카메라 없이 움직임과 재실 상태를 감지하고,\n장시간 무활동 시 대응 절차를 연결합니다.",
        Inches(0.78),
        Inches(3.9),
        Inches(6.3),
        Inches(0.8),
        size=15,
        color=RGBColor(189, 222, 212),
    )
    add_text(
        slide,
        "팀 올웨이즈(always)  |  동의대학교 응용소프트웨어공학과  |  2026",
        Inches(0.78),
        Inches(6.55),
        Inches(7.3),
        Inches(0.35),
        size=11,
        color=RGBColor(189, 222, 212),
    )
    for index, (label, value, color) in enumerate(
        [
            ("SENSE", "비접촉 감지", COLORS["mint"]),
            ("ANALYZE", "무활동 판정", RGBColor(250, 190, 88)),
            ("RESPOND", "안전 대응", RGBColor(244, 116, 128)),
        ]
    ):
        y = Inches(1.18 + index * 1.75)
        add_text(
            slide,
            f"0{index + 1}",
            Inches(9.55),
            y,
            Inches(0.55),
            Inches(0.35),
            size=12,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            label,
            Inches(10.15),
            y,
            Inches(1.8),
            Inches(0.32),
            size=12,
            color=COLORS["white"],
            bold=True,
        )
        add_text(
            slide,
            value,
            Inches(9.55),
            y + Inches(0.48),
            Inches(2.4),
            Inches(0.42),
            size=18,
            color=RGBColor(225, 244, 237),
            bold=True,
        )

    # 2. Problem and goal
    slide = prs.slides.add_slide(blank)
    add_title(slide, "01", "프로젝트 배경과 목표", "계획서가 정의한 사회적 문제와 핵심 설계 원칙")
    add_rect(
        slide,
        Inches(0.6),
        Inches(1.92),
        Inches(5.9),
        Inches(4.72),
        COLORS["pale"],
        line=COLORS["pale"],
    )
    add_text(
        slide,
        "기존 방식의 한계",
        Inches(0.95),
        Inches(2.2),
        Inches(2.8),
        Inches(0.4),
        size=21,
        color=COLORS["navy"],
        bold=True,
    )
    for index, (head, body) in enumerate(
        [
            ("카메라", "사생활 노출과 심리적 거부감"),
            ("웨어러블", "착용·충전 등 지속적 참여 필요"),
            ("인력 중심 확인", "24시간 상시 대응과 확장성 한계"),
        ]
    ):
        y = Inches(2.85 + index * 1.05)
        add_rect(
            slide,
            Inches(0.95),
            y,
            Inches(0.54),
            Inches(0.54),
            COLORS["white"],
            line=COLORS["line"],
        )
        add_text(
            slide,
            str(index + 1),
            Inches(0.95),
            y,
            Inches(0.54),
            Inches(0.54),
            size=14,
            color=COLORS["red"],
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            head,
            Inches(1.7),
            y - Inches(0.02),
            Inches(1.3),
            Inches(0.32),
            size=15,
            color=COLORS["ink"],
            bold=True,
        )
        add_text(
            slide,
            body,
            Inches(1.7),
            y + Inches(0.34),
            Inches(3.9),
            Inches(0.32),
            size=11.5,
            color=COLORS["muted"],
        )
    add_rect(
        slide,
        Inches(6.82),
        Inches(1.92),
        Inches(5.9),
        Inches(4.72),
        COLORS["navy"],
        line=COLORS["navy"],
    )
    add_text(
        slide,
        "LoneCare의 목표",
        Inches(7.18),
        Inches(2.2),
        Inches(3.3),
        Inches(0.4),
        size=21,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "카메라 없이\n24시간 생활 안전을 확인",
        Inches(7.18),
        Inches(2.85),
        Inches(4.8),
        Inches(1.1),
        size=27,
        color=COLORS["mint"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "PIR·압력·레이더 센서로 움직임과 재실 상태 수집",
            "12시간 무활동 기준으로 주의·위험 상태 자동 판정",
            "보호자·담당자·현장 확인으로 이어지는 대응 절차",
            "다수 가구를 한 화면에서 관리하는 웹 관제 시스템",
        ],
        Inches(7.15),
        Inches(4.25),
        Inches(4.9),
        Inches(1.75),
        size=13.5,
        color=RGBColor(225, 244, 237),
        gap=10,
    )
    add_footer(slide, 2, "근거: 임베디드시스템 프로젝트 계획서")

    # 3. Plan vs implementation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "02", "계획 대비 현재 구현", "초기 목표를 유지하면서 운영 기능을 확장")
    headers = [("계획서 목표", COLORS["blue_pale"], COLORS["blue"]), ("현재 구현", COLORS["pale"], COLORS["green"])]
    for index, (header, fill, accent) in enumerate(headers):
        x = Inches(0.62 + index * 6.15)
        add_rect(slide, x, Inches(1.88), Inches(5.85), Inches(4.75), fill, line=fill)
        add_text(
            slide,
            header,
            x + Inches(0.3),
            Inches(2.15),
            Inches(2.3),
            Inches(0.4),
            size=20,
            color=accent,
            bold=True,
        )
    add_bullets(
        slide,
        [
            "PIR + 압력 센서 기반 활동 감지",
            "ESP32/아두이노 → 라즈베리 파이 → 웹 서버",
            "12시간 무활동 시 이상 징후 판정",
            "보호자·관리 기관에 경고 알림",
            "실시간 대시보드와 다가구 관제",
        ],
        Inches(0.95),
        Inches(2.85),
        Inches(5.0),
        Inches(2.85),
        size=16,
        gap=12,
    )
    add_bullets(
        slide,
        [
            "LD2410 레이더까지 확장한 재실·움직임 감지",
            "FastAPI + SQLite + Docker 기반 단일 서버 구성",
            "상태·알림·장치·보고서·시뮬레이션 API",
            "보호자 → 관리자 → 현장 방문 단계별 워크플로",
            "필터·활동 그래프·안전 확인 이력 관리 UI",
        ],
        Inches(7.1),
        Inches(2.85),
        Inches(5.0),
        Inches(2.85),
        size=16,
        gap=12,
    )
    add_rect(
        slide,
        Inches(4.68),
        Inches(5.95),
        Inches(3.95),
        Inches(0.48),
        COLORS["white"],
        line=COLORS["line"],
    )
    add_text(
        slide,
        "핵심 방향은 유지, 관제 운영 기능은 확대",
        Inches(4.68),
        Inches(5.95),
        Inches(3.95),
        Inches(0.48),
        size=12,
        color=COLORS["navy"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, 3)

    # 4. Architecture
    slide = prs.slides.add_slide(blank)
    add_title(slide, "03", "시스템 아키텍처", "센서 노드부터 관제 화면까지 이어지는 데이터 경로")
    positions = [
        (Inches(0.65), "센서 노드", "PIR · 압력 · LD2410\n생활 활동 감지", COLORS["blue"], "1"),
        (Inches(3.2), "ESP32", "C/C++ 펌웨어\nWi-Fi 설정 · JSON 전송", COLORS["green"], "2"),
        (Inches(5.75), "FastAPI 서버", "REST API · 상태 판정\n알림·보고서 서비스", COLORS["amber"], "3"),
        (Inches(8.3), "SQLite", "장치 · 센서 로그\n알림 · 처리 이력", COLORS["red"], "4"),
        (Inches(10.85), "웹 관제", "대시보드 · 장치 관리\n안전 대응 워크플로", COLORS["navy"], "5"),
    ]
    for x, title, detail, color, number in positions:
        add_arch_box(slide, title, detail, x, Inches(2.15), Inches(1.85), color, number)
    for index in range(len(positions) - 1):
        connect(
            slide,
            positions[index][0] + Inches(1.87),
            Inches(2.78),
            positions[index + 1][0] - Inches(0.03),
            Inches(2.78),
        )
    add_rect(
        slide,
        Inches(0.72),
        Inches(4.25),
        Inches(11.95),
        Inches(1.65),
        COLORS["gray"],
        line=COLORS["line"],
    )
    add_text(
        slide,
        "배포 구조",
        Inches(1.0),
        Inches(4.55),
        Inches(1.3),
        Inches(0.35),
        size=15,
        color=COLORS["navy"],
        bold=True,
    )
    add_rich_text(
        slide,
        [
            ("Docker Compose", True, COLORS["green"]),
            ("로 Windows 개발 환경과 Raspberry Pi ARM64 운영 환경을 공통 구성", False, COLORS["ink"]),
        ],
        Inches(2.25),
        Inches(4.47),
        Inches(9.6),
        Inches(0.45),
        size=14,
    )
    add_text(
        slide,
        "Named volume으로 SQLite DB를 보존하고, 헬스체크와 자동 재시작으로 운영 안정성을 확보합니다.",
        Inches(2.25),
        Inches(5.05),
        Inches(9.6),
        Inches(0.42),
        size=12,
        color=COLORS["muted"],
    )
    add_footer(slide, 4)

    # 5. Decision algorithm
    slide = prs.slides.add_slide(blank)
    add_title(slide, "04", "센서 데이터 처리와 위험 판정", "수신 패킷을 활동 여부로 변환하고 상태·알림을 갱신")
    steps = [
        ("01", "데이터 수신", "PIR · 레이더 · 압력\n배터리 · Wi-Fi"),
        ("02", "활동 판단", "PIR 움직임 또는\n압력 변화량 ≥ 임계값"),
        ("03", "시간 계산", "마지막 활동 시각부터\n무활동 시간 계산"),
        ("04", "상태 판정", "normal → warning → danger"),
        ("05", "알림 처리", "위험 알림 1건 유지\n활동 재감지 시 자동 해제"),
    ]
    for index, (num, head, detail) in enumerate(steps):
        x = Inches(0.7 + index * 2.52)
        add_rect(
            slide,
            x,
            Inches(2.05),
            Inches(2.18),
            Inches(1.7),
            COLORS["white"],
            line=COLORS["line"],
        )
        add_text(
            slide,
            num,
            x + Inches(0.2),
            Inches(2.28),
            Inches(0.5),
            Inches(0.3),
            size=11,
            color=COLORS["green"],
            bold=True,
        )
        add_text(
            slide,
            head,
            x + Inches(0.2),
            Inches(2.65),
            Inches(1.75),
            Inches(0.35),
            size=15,
            color=COLORS["navy"],
            bold=True,
        )
        add_text(
            slide,
            detail,
            x + Inches(0.2),
            Inches(3.1),
            Inches(1.78),
            Inches(0.48),
            size=10,
            color=COLORS["muted"],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if index < 4:
            connect(
                slide,
                x + Inches(2.2),
                Inches(2.9),
                x + Inches(2.48),
                Inches(2.9),
            )
    add_status_card(
        slide,
        "안전 NORMAL",
        "< 6h",
        "최근 활동 확인",
        Inches(1.05),
        Inches(4.45),
        COLORS["green"],
        RGBColor(220, 244, 235),
    )
    add_status_card(
        slide,
        "주의 WARNING",
        "6~12h",
        "선제 확인 필요",
        Inches(5.3),
        Inches(4.45),
        COLORS["amber"],
        COLORS["amber_pale"],
    )
    add_status_card(
        slide,
        "위험 DANGER",
        "≥ 12h",
        "즉시 안전 대응",
        Inches(9.55),
        Inches(4.45),
        COLORS["red"],
        COLORS["red_pale"],
    )
    add_text(
        slide,
        "※ 운영 기준은 환경 변수로 변경 가능하며, 현재 코드는 warning과 danger 임계값을 독립 설정합니다.",
        Inches(0.8),
        Inches(6.5),
        Inches(11.8),
        Inches(0.32),
        size=10,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 5, "근거: app/routers/sensor.py · app/services/status_service.py")

    # 6. Code structure
    slide = prs.slides.add_slide(blank)
    add_title(slide, "05", "코드 구조와 핵심 모듈", "기능별 Router와 Service를 분리한 FastAPI 애플리케이션")
    modules = [
        ("routers", "sensor · status · alerts\n devices · reports · simulation", COLORS["blue"]),
        ("services", "상태 계산\n위험 알림 생성·해제", COLORS["green"]),
        ("models", "Device · SensorLog\nAlert · AlertActionLog", COLORS["amber"]),
        ("static", "대시보드 · 장치 관리\n안전 확인 로그", COLORS["red"]),
    ]
    for index, (head, detail, color) in enumerate(modules):
        x = Inches(0.7 + index * 3.12)
        add_rect(
            slide,
            x,
            Inches(2.0),
            Inches(2.75),
            Inches(1.55),
            COLORS["white"],
            line=COLORS["line"],
        )
        add_rect(slide, x, Inches(2.0), Inches(2.75), Inches(0.12), color, radius=False)
        add_text(
            slide,
            head,
            x + Inches(0.25),
            Inches(2.32),
            Inches(2.1),
            Inches(0.35),
            size=18,
            color=COLORS["navy"],
            bold=True,
        )
        add_text(
            slide,
            detail,
            x + Inches(0.25),
            Inches(2.85),
            Inches(2.2),
            Inches(0.5),
            size=11.5,
            color=COLORS["muted"],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    add_rect(
        slide,
        Inches(0.72),
        Inches(4.05),
        Inches(6.0),
        Inches(2.2),
        COLORS["pale"],
        line=COLORS["pale"],
    )
    add_text(
        slide,
        "구현 강점",
        Inches(1.02),
        Inches(4.35),
        Inches(2.0),
        Inches(0.35),
        size=18,
        color=COLORS["green"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "새 device_id 자동 등록과 부분 센서 업데이트",
            "열린 danger 알림 중복 생성 방지",
            "활동 재감지 시 상태 정상화와 알림 자동 해제",
            "처리 단계·담당자 연락망·완료 이력까지 데이터화",
        ],
        Inches(1.0),
        Inches(4.85),
        Inches(5.3),
        Inches(1.2),
        size=12.5,
        gap=7,
    )
    add_rect(
        slide,
        Inches(6.95),
        Inches(4.05),
        Inches(5.65),
        Inches(2.2),
        COLORS["gray"],
        line=COLORS["line"],
    )
    add_text(
        slide,
        "현재 범위와 개선점",
        Inches(7.25),
        Inches(4.35),
        Inches(2.5),
        Inches(0.35),
        size=18,
        color=COLORS["red"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "실제 SMS·Push 발송 모듈은 아직 미연동",
            "사용자 인증·권한 관리 기능 보강 필요",
            "계획서의 Raspberry Pi 게이트웨이는 현재 저장소에 미포함",
            "SQLite 단일 서버 구조는 대규모 운영 전 DB 전환 필요",
        ],
        Inches(7.22),
        Inches(4.85),
        Inches(4.9),
        Inches(1.2),
        size=12.5,
        gap=7,
    )
    add_footer(slide, 6, "근거: app/ 전체 코드 분석")

    # 7. Dashboard
    slide = prs.slides.add_slide(blank)
    add_title(slide, "06", "통합 대시보드", "위험 대상 우선 노출과 실시간 센서·알림 확인")
    add_picture_cover(
        slide,
        crops["dashboard_top"],
        Inches(0.58),
        Inches(1.82),
        Inches(8.35),
        Inches(4.95),
    )
    callouts = [
        ("51대", "전체 등록 장치", COLORS["navy"]),
        ("26 / 15 / 10", "안전 · 주의 · 위험", COLORS["green"]),
        ("위험 우선", "상위 6개 장치 카드", COLORS["red"]),
        ("5초", "자동 갱신 주기", COLORS["blue"]),
    ]
    for index, (value, label, color) in enumerate(callouts):
        y = Inches(1.92 + index * 1.12)
        add_rect(
            slide,
            Inches(9.25),
            y,
            Inches(3.42),
            Inches(0.88),
            COLORS["white"],
            line=COLORS["line"],
        )
        add_text(
            slide,
            value,
            Inches(9.52),
            y + Inches(0.14),
            Inches(1.5),
            Inches(0.35),
            size=18,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            label,
            Inches(10.85),
            y + Inches(0.23),
            Inches(1.52),
            Inches(0.26),
            size=10.5,
            color=COLORS["muted"],
            align=PP_ALIGN.RIGHT,
        )
    add_text(
        slide,
        "발표 시연 포인트",
        Inches(9.28),
        Inches(6.25),
        Inches(1.7),
        Inches(0.3),
        size=12,
        color=COLORS["green"],
        bold=True,
    )
    add_text(
        slide,
        "상태 변화 → 위험 알림 → 안전 확인 완료를 한 화면에서 확인",
        Inches(10.65),
        Inches(6.18),
        Inches(1.9),
        Inches(0.48),
        size=10.5,
        color=COLORS["ink"],
        align=PP_ALIGN.RIGHT,
    )
    add_footer(slide, 7, "직접 촬영: http://localhost:8000")

    # 8. Device management
    slide = prs.slides.add_slide(blank)
    add_title(slide, "07", "장치 관리와 이상 대상 분류", "다수 가구 운영을 위한 검색·필터·정렬·처리 기능")
    add_picture_cover(
        slide,
        crops["devices"],
        Inches(0.58),
        Inches(1.82),
        Inches(8.7),
        Inches(4.98),
    )
    add_rect(
        slide,
        Inches(9.55),
        Inches(1.95),
        Inches(3.1),
        Inches(4.65),
        COLORS["pale"],
        line=COLORS["pale"],
    )
    add_text(
        slide,
        "운영자 기능",
        Inches(9.85),
        Inches(2.25),
        Inches(1.7),
        Inches(0.35),
        size=19,
        color=COLORS["navy"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "상태·연결·위치·센서 조건별 필터",
            "위험 우선 / 최근 수신 / 무활동 정렬",
            "배터리·Wi-Fi·침대·사람 존재 상태 확인",
            "24시간 활동 그래프와 안전 확인 사유 조회",
            "보호자·담당자 연락처와 우선순위 관리",
        ],
        Inches(9.8),
        Inches(2.9),
        Inches(2.45),
        Inches(2.7),
        size=13,
        gap=10,
    )
    add_rect(
        slide,
        Inches(9.82),
        Inches(5.72),
        Inches(2.55),
        Inches(0.55),
        COLORS["navy"],
        line=COLORS["navy"],
    )
    add_text(
        slide,
        "복지 담당자의 관제 업무를 화면 중심으로 표준화",
        Inches(9.9),
        Inches(5.75),
        Inches(2.38),
        Inches(0.48),
        size=10,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, 8, "직접 촬영: http://localhost:8000/devices")

    # 9. Workflow
    slide = prs.slides.add_slide(blank)
    add_title(slide, "08", "위험 알림 대응 워크플로", "감지에서 현장 확인까지 처리 단계를 기록")
    stages = [
        ("위험 감지", COLORS["red"]),
        ("보호자 알림", RGBColor(227, 105, 72)),
        ("응답 대기", COLORS["amber"]),
        ("관리자 확인", COLORS["blue"]),
        ("현장 방문", COLORS["green"]),
        ("안전 확인", COLORS["navy"]),
    ]
    for index, (name, color) in enumerate(stages):
        x = Inches(0.62 + index * 2.08)
        add_rect(
            slide,
            x,
            Inches(1.92),
            Inches(1.72),
            Inches(0.72),
            color,
            line=color,
        )
        add_text(
            slide,
            name,
            x,
            Inches(1.92),
            Inches(1.72),
            Inches(0.72),
            size=12,
            color=COLORS["white"],
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if index < 5:
            connect(
                slide,
                x + Inches(1.72),
                Inches(2.28),
                x + Inches(2.02),
                Inches(2.28),
                COLORS["muted"],
            )
    add_picture_cover(
        slide,
        crops["resolutions"],
        Inches(0.62),
        Inches(2.95),
        Inches(7.7),
        Inches(3.65),
    )
    add_rect(
        slide,
        Inches(8.62),
        Inches(2.95),
        Inches(4.05),
        Inches(3.65),
        COLORS["gray"],
        line=COLORS["line"],
    )
    add_text(
        slide,
        "처리 이력에 남는 정보",
        Inches(8.95),
        Inches(3.25),
        Inches(2.8),
        Inches(0.35),
        size=18,
        color=COLORS["navy"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "위험이 발생한 원인과 시각",
            "보호자·담당자 전달 및 단계 변경 로그",
            "현장 방문·전화·센서 점검 등 확인 방식",
            "최종 처리 단계와 상세 메모",
            "새로고침 후에도 유지되는 공식 기록",
        ],
        Inches(8.92),
        Inches(3.9),
        Inches(3.18),
        Inches(2.25),
        size=13,
        gap=10,
    )
    add_footer(slide, 9, "직접 촬영: http://localhost:8000/resolutions")

    # 10. API and operation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "09", "REST API와 운영 배포", "센서·관제 화면·운영 도구가 동일한 API를 사용")
    add_picture_cover(
        slide,
        crops["docs"],
        Inches(0.58),
        Inches(1.82),
        Inches(7.55),
        Inches(4.95),
    )
    add_rect(
        slide,
        Inches(8.42),
        Inches(1.88),
        Inches(4.25),
        Inches(2.05),
        COLORS["pale"],
        line=COLORS["pale"],
    )
    add_text(
        slide,
        "주요 API",
        Inches(8.75),
        Inches(2.15),
        Inches(1.5),
        Inches(0.35),
        size=18,
        color=COLORS["navy"],
        bold=True,
    )
    add_text(
        slide,
        "POST  /api/sensor-data\nGET   /api/status · /api/activity\nGET   /api/alerts · /api/reports\nPOST  /api/alerts/{id}/workflow/action",
        Inches(8.75),
        Inches(2.65),
        Inches(3.45),
        Inches(1.05),
        size=11.5,
        color=COLORS["ink"],
        font="Cascadia Mono",
    )
    add_rect(
        slide,
        Inches(8.42),
        Inches(4.18),
        Inches(4.25),
        Inches(2.35),
        COLORS["gray"],
        line=COLORS["line"],
    )
    add_text(
        slide,
        "운영 구성",
        Inches(8.75),
        Inches(4.45),
        Inches(1.5),
        Inches(0.35),
        size=18,
        color=COLORS["navy"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "Docker Compose 단일 명령 실행",
            "Windows 개발: Uvicorn reload",
            "Raspberry Pi: restart unless-stopped",
            "Named volume DB 영구 저장",
            "/health 기반 컨테이너 헬스체크",
        ],
        Inches(8.72),
        Inches(4.95),
        Inches(3.4),
        Inches(1.35),
        size=12.5,
        gap=7,
    )
    add_footer(slide, 10, "직접 촬영: http://localhost:8000/docs")

    # 11. Test results
    slide = prs.slides.add_slide(blank)
    add_title(slide, "10", "실행 검증과 시연 데이터", "Docker 환경에서 실제 화면과 API 응답을 확인")
    metrics = [
        ("51", "등록 장치", COLORS["navy"], COLORS["pale"]),
        ("26", "안전", COLORS["green"], RGBColor(220, 244, 235)),
        ("15", "주의", COLORS["amber"], COLORS["amber_pale"]),
        ("10", "위험", COLORS["red"], COLORS["red_pale"]),
        ("10", "미해결 알림", COLORS["red"], COLORS["red_pale"]),
    ]
    for index, (value, label, color, pale) in enumerate(metrics):
        x = Inches(0.62 + index * 2.52)
        add_rect(
            slide,
            x,
            Inches(1.95),
            Inches(2.18),
            Inches(1.35),
            COLORS["white"],
            line=COLORS["line"],
        )
        add_rect(
            slide,
            x + Inches(1.65),
            Inches(2.55),
            Inches(0.35),
            Inches(0.35),
            pale,
        )
        add_text(
            slide,
            value,
            x + Inches(0.22),
            Inches(2.18),
            Inches(1.2),
            Inches(0.5),
            size=29,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            label,
            x + Inches(0.22),
            Inches(2.75),
            Inches(1.5),
            Inches(0.25),
            size=10.5,
            color=COLORS["muted"],
        )
    add_rect(
        slide,
        Inches(0.62),
        Inches(3.65),
        Inches(5.95),
        Inches(2.65),
        COLORS["pale"],
        line=COLORS["pale"],
    )
    add_text(
        slide,
        "검증한 동작",
        Inches(0.95),
        Inches(3.95),
        Inches(1.8),
        Inches(0.35),
        size=18,
        color=COLORS["green"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "컨테이너 health 상태 정상",
            "샘플 데이터 50대와 기존 센서 장치 동시 표시",
            "danger 장치별 열린 알림 1건 유지",
            "장치·로그·알림 데이터 volume 영구 저장",
        ],
        Inches(0.92),
        Inches(4.48),
        Inches(5.1),
        Inches(1.45),
        size=13.5,
        gap=9,
    )
    add_rect(
        slide,
        Inches(6.82),
        Inches(3.65),
        Inches(5.85),
        Inches(2.65),
        COLORS["gray"],
        line=COLORS["line"],
    )
    add_text(
        slide,
        "발표 시연 순서",
        Inches(7.15),
        Inches(3.95),
        Inches(1.8),
        Inches(0.35),
        size=18,
        color=COLORS["navy"],
        bold=True,
    )
    add_bullets(
        slide,
        [
            "1. 센서 데이터 전송 또는 시뮬레이션",
            "2. normal → warning → danger 상태 변화",
            "3. 위험 알림과 대응 단계 확인",
            "4. 안전 확인 완료 후 정상 복구·이력 저장",
        ],
        Inches(7.12),
        Inches(4.48),
        Inches(5.0),
        Inches(1.45),
        size=13.5,
        gap=9,
    )
    add_footer(slide, 11, "검증일: 2026-06-14 · Docker Desktop 로컬 실행")

    # 12. Conclusion
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["navy"], radius=False)
    add_text(
        slide,
        "11  CONCLUSION",
        Inches(0.72),
        Inches(0.68),
        Inches(2.5),
        Inches(0.3),
        size=11,
        color=COLORS["mint"],
        bold=True,
    )
    add_text(
        slide,
        "프라이버시를 지키면서\n위험 신호를 놓치지 않는 생활 안전망",
        Inches(0.72),
        Inches(1.35),
        Inches(7.6),
        Inches(1.45),
        size=31,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "IoT LoneCare는 센서 데이터 수집을 넘어\n감지 → 판정 → 대응 → 기록의 운영 흐름을 구현했습니다.",
        Inches(0.75),
        Inches(3.1),
        Inches(6.8),
        Inches(0.85),
        size=17,
        color=RGBColor(205, 231, 223),
    )
    roadmap = [
        ("NEXT 01", "실제 알림 연동", "SMS · Push · 기관 메시징"),
        ("NEXT 02", "보안 강화", "인증 · 권한 · 전송 암호화"),
        ("NEXT 03", "판정 고도화", "개인별 기준 · 오탐 감소"),
        ("NEXT 04", "현장 확장", "Raspberry Pi 게이트웨이 · 다기관 운영"),
    ]
    for index, (tag, head, detail) in enumerate(roadmap):
        x = Inches(7.85 + (index % 2) * 2.55)
        y = Inches(1.25 + (index // 2) * 2.35)
        add_rect(
            slide,
            x,
            y,
            Inches(2.25),
            Inches(1.75),
            RGBColor(26, 83, 70),
            line=RGBColor(47, 112, 95),
        )
        add_text(
            slide,
            tag,
            x + Inches(0.22),
            y + Inches(0.2),
            Inches(1.5),
            Inches(0.25),
            size=9,
            color=COLORS["mint"],
            bold=True,
        )
        add_text(
            slide,
            head,
            x + Inches(0.22),
            y + Inches(0.58),
            Inches(1.8),
            Inches(0.35),
            size=16,
            color=COLORS["white"],
            bold=True,
        )
        add_text(
            slide,
            detail,
            x + Inches(0.22),
            y + Inches(1.08),
            Inches(1.8),
            Inches(0.42),
            size=10,
            color=RGBColor(188, 220, 211),
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "감사합니다",
        Inches(0.75),
        Inches(6.45),
        Inches(2.8),
        Inches(0.55),
        size=27,
        color=COLORS["mint"],
        bold=True,
    )
    add_text(
        slide,
        "Q & A",
        Inches(10.45),
        Inches(6.48),
        Inches(1.8),
        Inches(0.5),
        size=23,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(create_presentation())
